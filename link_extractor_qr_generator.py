#!/usr/bin/env python3
"""
Link Extractor and QR Code Generator

This script extracts links from various text-based files and generates QR codes,
organizing them in a structured directory format with multiple output options.
"""

import os
import re
import yaml
import json
import qrcode
import argparse
import shutil
from pathlib import Path
from urllib.parse import urlparse
import uuid
from typing import List, Dict, Optional, Set, Any, Union

# Import optional modules with error handling
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from fpdf import FPDF
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    import PyPDF2
    PYPDF2_AVAILABLE = True
except ImportError:
    PYPDF2_AVAILABLE = False

try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

try:
    from bs4 import BeautifulSoup
    BEAUTIFULSOUP_AVAILABLE = True
except ImportError:
    BEAUTIFULSOUP_AVAILABLE = False


class LinkExtractor:
    """Extract links from various file types."""
    
    def __init__(self, config: Dict[str, Any]):
        """
        Initialize the LinkExtractor with configuration.
        
        Args:
            config: Configuration dictionary
        """
        self.config = config
        self.url_pattern = re.compile(r'http[s]?://(?:[a-zA-Z]|[0-9]|[$-_@.&+]|[!*\\(\\),]|(?:%[0-9a-fA-F][0-9a-fA-F]))+')
    
    def extract_from_file(self, file_path: str) -> Set[str]:
        """
        Extract links from a file based on its extension.
        
        Args:
            file_path: Path to the file
            
        Returns:
            Set of unique URLs found in the file
        """
        file_extension = os.path.splitext(file_path)[1].lower()
        
        extractors = {
            '.txt': self._extract_from_text,
            '.md': self._extract_from_text,
            '.csv': self._extract_from_csv,
            '.pdf': self._extract_from_pdf,
            '.docx': self._extract_from_docx,
            '.doc': self._extract_from_text,  # Fallback for .doc files
            '.html': self._extract_from_html,
            '.htm': self._extract_from_html,
            '.yaml': self._extract_from_yaml,
            '.yml': self._extract_from_yaml,
            '.json': self._extract_from_json
        }
        
        extractor = extractors.get(file_extension, self._extract_from_text)
        try:
            return extractor(file_path)
        except Exception as e:
            print(f"Error extracting links from {file_path}: {e}")
            return set()
    
    def _extract_from_text(self, file_path: str) -> Set[str]:
        """Extract links from text-based files."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                content = file.read()
            return self._find_urls(content)
        except Exception as e:
            print(f"Error reading text file {file_path}: {e}")
            return set()
    
    def _extract_from_csv(self, file_path: str) -> Set[str]:
        """Extract links from CSV files."""
        if not PANDAS_AVAILABLE:
            print("Pandas is not installed. Falling back to text extraction for CSV.")
            return self._extract_from_text(file_path)
        
        try:
            df = pd.read_csv(file_path)
            urls = set()
            for column in df.columns:
                for value in df[column].astype(str):
                    urls.update(self._find_urls(value))
            return urls
        except Exception as e:
            print(f"Error processing CSV file {file_path}: {e}")
            return self._extract_from_text(file_path)
    
    def _extract_from_pdf(self, file_path: str) -> Set[str]:
        """Extract links from PDF files."""
        if not PYPDF2_AVAILABLE:
            print("PyPDF2 is not installed. Falling back to text extraction for PDF.")
            return self._extract_from_text(file_path)
        
        try:
            pdf_file = open(file_path, 'rb')
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            urls = set()
            for page in pdf_reader.pages:
                text = page.extract_text()
                if text:
                    urls.update(self._find_urls(text))
            
            pdf_file.close()
            return urls
        except Exception as e:
            print(f"Error processing PDF file {file_path}: {e}")
            return set()
    
    def _extract_from_docx(self, file_path: str) -> Set[str]:
        """Extract links from DOCX files."""
        if not DOCX_AVAILABLE:
            print("python-docx is not installed. Falling back to text extraction for DOCX.")
            return self._extract_from_text(file_path)
        
        try:
            doc = docx.Document(file_path)
            content = "\n".join([para.text for para in doc.paragraphs])
            return self._find_urls(content)
        except Exception as e:
            print(f"Error processing DOCX file {file_path}: {e}")
            return set()
    
    def _extract_from_html(self, file_path: str) -> Set[str]:
        """Extract links from HTML files."""
        if not BEAUTIFULSOUP_AVAILABLE:
            print("BeautifulSoup is not installed. Falling back to text extraction for HTML.")
            return self._extract_from_text(file_path)
        
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                content = file.read()
            
            soup = BeautifulSoup(content, 'html.parser')
            
            # Extract URLs from href attributes
            urls = set()
            for a_tag in soup.find_all('a', href=True):
                href = a_tag['href']
                if href.startswith('http'):
                    urls.add(href)
            
            # Also find URLs in the text using regex for any missed links
            urls.update(self._find_urls(content))
            
            return urls
        except Exception as e:
            print(f"Error processing HTML file {file_path}: {e}")
            return self._extract_from_text(file_path)
    
    def _extract_from_yaml(self, file_path: str) -> Set[str]:
        """Extract links from YAML files."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                content = file.read()
                
            # First, try to find URLs in the raw text
            urls = self._find_urls(content)
            
            # Then try to parse YAML and search for URLs in values
            try:
                yaml_data = yaml.safe_load(content)
                if yaml_data:
                    self._find_urls_in_structure(yaml_data, urls)
            except Exception as yaml_e:
                print(f"Error parsing YAML structure in {file_path}: {yaml_e}")
                # Already found URLs from raw text, so continue
                
            return urls
        except Exception as e:
            print(f"Error processing YAML file {file_path}: {e}")
            return set()
    
    def _extract_from_json(self, file_path: str) -> Set[str]:
        """Extract links from JSON files."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                content = file.read()
                
            # First, try to find URLs in the raw text
            urls = self._find_urls(content)
            
            # Then try to parse JSON and search for URLs in values
            try:
                json_data = json.loads(content)
                self._find_urls_in_structure(json_data, urls)
            except Exception as json_e:
                print(f"Error parsing JSON structure in {file_path}: {json_e}")
                # Already found URLs from raw text, so continue
                
            return urls
        except Exception as e:
            print(f"Error processing JSON file {file_path}: {e}")
            return set()
    
    def _find_urls_in_structure(self, data: Any, urls: Set[str]) -> None:
        """Recursively search for URLs in nested data structures (dict, list, etc.)."""
        if isinstance(data, dict):
            for value in data.values():
                self._find_urls_in_structure(value, urls)
        elif isinstance(data, list):
            for item in data:
                self._find_urls_in_structure(item, urls)
        elif isinstance(data, str):
            urls.update(self._find_urls(data))
    
    def _find_urls(self, text: str) -> Set[str]:
        """Find all URLs in a text string."""
        return set(self.url_pattern.findall(text))


class QRCodeGenerator:
    """Generate QR codes from URLs."""
    
    def __init__(self, config: Dict[str, Any]):
        """
        Initialize the QRCodeGenerator with configuration.
        
        Args:
            config: Configuration dictionary
        """
        self.config = config
        self.error_correction = getattr(qrcode.constants, 
                                      config.get('qr_error_correction', 'ERROR_CORRECT_L'))
        self.box_size = config.get('qr_box_size', 10)
        self.border = config.get('qr_border', 4)
    
    def generate_qr_code(self, url: str, output_path: str, filename: str = "qrcode") -> Dict[str, str]:
        """
        Generate QR code for a URL and save as image and binary file.
        
        Args:
            url: URL to encode in QR code
            output_path: Directory to save the QR code files
            filename: Base filename for output files
            
        Returns:
            Dictionary containing paths to generated files
        """
        os.makedirs(output_path, exist_ok=True)
        
        # Create QR code
        qr = qrcode.QRCode(
            version=None,  # Auto-determine
            error_correction=self.error_correction,
            box_size=self.box_size,
            border=self.border,
        )
        qr.add_data(url)
        qr.make(fit=True)
        
        img = qr.make_image(fill_color="black", back_color="white")
        
        # Save image
        image_path = os.path.join(output_path, f"{filename}.png")
        img.save(image_path)
        
        # Save binary data
        binary_path = os.path.join(output_path, f"{filename}.bin")
        self._save_binary_data(qr, url, binary_path)
        
        # Save URL to markdown file
        url_path = os.path.join(output_path, "url.md")
        with open(url_path, 'w', encoding='utf-8') as f:
            f.write(f"# QR Code Link\n\n{url}\n")
        
        return {
            'image': image_path,
            'binary': binary_path,
            'url': url_path
        }
    
    def _save_binary_data(self, qr: qrcode.QRCode, url: str, file_path: str) -> None:
        """Save QR code data as binary file."""
        data = qr.get_matrix()
        binary_data = ""
        for row in data:
            for cell in row:
                binary_data += "1" if cell else "0"
        
        additional_info = f"Link: {url}\nVersion: {qr.version}\n"
        data_file_content = additional_info.encode('utf-8') + binary_data.encode('utf-8')
        
        with open(file_path, "wb") as f:
            f.write(data_file_content)


class OutputGenerator:
    """Generate various output formats from QR codes."""
    
    def __init__(self, config: Dict[str, Any]):
        """
        Initialize the OutputGenerator with configuration.
        
        Args:
            config: Configuration dictionary
        """
        self.config = config
    
    def create_powerpoint(self, qr_data: List[Dict[str, Any]], output_path: str) -> Optional[str]:
        """
        Create PowerPoint presentation with QR codes and links.
        
        Args:
            qr_data: List of dictionaries containing QR code data
            output_path: Path to save the PowerPoint file
            
        Returns:
            Path to the created PowerPoint file or None if failed
        """
        if not PPTX_AVAILABLE:
            print("python-pptx is not installed. PowerPoint creation skipped.")
            return None
        
        if not qr_data:
            print("No QR codes to include in PowerPoint.")
            return None
        
        try:
            prs = Presentation()
            
            for data in qr_data:
                slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
                
                # Add title (URL name/description)
                title = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.5), Inches(9), Inches(1)
                )
                title_frame = title.text_frame
                p = title_frame.add_paragraph()
                p.text = data.get('title', data['url'])
                p.font.size = Pt(24)
                
                # Add QR code image
                img_path = data['files']['image']
                slide.shapes.add_picture(
                    img_path, Inches(2), Inches(1.5), width=Inches(4)
                )
                
                # Add URL as clickable link
                link_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(6), Inches(9), Inches(1)
                )
                link_frame = link_box.text_frame
                p = link_frame.add_paragraph()
                r = p.add_run()
                r.text = data['url']
                
                # Try to make it a hyperlink if possible
                try:
                    r.hyperlink.address = data['url']
                except:
                    pass  # Some versions of python-pptx don't support this
            
            pptx_path = os.path.join(output_path, "qr_codes.pptx")
            prs.save(pptx_path)
            return pptx_path
        
        except Exception as e:
            print(f"Error creating PowerPoint: {e}")
            return None
    
    def create_pdf(self, qr_data: List[Dict[str, Any]], output_path: str) -> Optional[str]:
        """
        Create PDF with QR codes and links.
        
        Args:
            qr_data: List of dictionaries containing QR code data
            output_path: Path to save the PDF file
            
        Returns:
            Path to the created PDF file or None if failed
        """
        if not PDF_AVAILABLE:
            print("fpdf is not installed. PDF creation skipped.")
            return None
        
        if not qr_data:
            print("No QR codes to include in PDF.")
            return None
        
        try:
            pdf = FPDF()
            
            for data in qr_data:
                pdf.add_page()
                
                # Add title
                pdf.set_font("Arial", "B", 16)
                title = data.get('title', data['url'])
                pdf.cell(0, 10, title, ln=True, align='C')
                
                # Add QR code image
                img_path = data['files']['image']
                pdf.image(img_path, x=70, y=50, w=70)
                
                # Add URL as text
                pdf.set_font("Arial", "", 12)
                pdf.set_y(130)  # Position below the QR code
                pdf.cell(0, 10, data['url'], ln=True, align='C', link=data['url'])
            
            pdf_path = os.path.join(output_path, "qr_codes.pdf")
            pdf.output(pdf_path)
            return pdf_path
        
        except Exception as e:
            print(f"Error creating PDF: {e}")
            return None


class LinkProcessorApp:
    """Main application class for processing links and generating outputs."""
    
    def __init__(self, config_path: Optional[str] = None, config: Optional[Dict[str, Any]] = None):
        """
        Initialize the application with optional configuration file or dictionary.
        
        Args:
            config_path: Path to configuration YAML file
            config: Configuration dictionary (overrides config_path if provided)
        """
        if config:
            self.config = config
        else:
            self.config = self._load_config(config_path)
        
        self.link_extractor = LinkExtractor(self.config)
        self.qr_generator = QRCodeGenerator(self.config)
        self.output_generator = OutputGenerator(self.config)
    
    def _load_config(self, config_path: Optional[str]) -> Dict[str, Any]:
        """Load configuration from YAML file or use defaults."""
        default_config = {
            'output_pptx': True,
            'output_pdf': True,
            'output_raw': True,
            'qr_error_correction': 'ERROR_CORRECT_L',
            'qr_box_size': 10,
            'qr_border': 4,
            'generate_descriptive_titles': True
        }
        
        if config_path and os.path.exists(config_path):
            try:
                with open(config_path, 'r', encoding='utf-8') as f:
                    user_config = yaml.safe_load(f)
                if user_config:
                    default_config.update(user_config)
            except Exception as e:
                print(f"Error loading config from {config_path}: {e}")
                print("Using default configuration.")
        
        return default_config
    
    def sanitize_filename(self, url: str) -> str:
        """
        Convert URL to a safe directory name.
        
        Args:
            url: URL to sanitize
            
        Returns:
            Sanitized directory name
        """
        # Extract domain for the start of the name
        parsed_url = urlparse(url)
        domain = parsed_url.netloc
        
        # Remove common prefixes
        path = parsed_url.path
        path = path.strip('/')
        path = path.replace('www.', '')
        
        # Combine parts
        name = f"{domain}_{path}"
        
        # Replace invalid characters
        invalid_chars = r'[<>:"/\\|?*]'
        name = re.sub(invalid_chars, '_', name)
        
        # Truncate if too long
        if len(name) > 100:
            name = name[:97] + "..."
        
        # Ensure uniqueness if the name is empty or too generic
        if not name or name == '_':
            name = f"url_{uuid.uuid4().hex[:8]}"
            
        return name
    
    def generate_title(self, url: str) -> str:
        """
        Generate a descriptive title for a URL.
        
        Args:
            url: URL to generate title for
            
        Returns:
            Generated title
        """
        if not self.config.get('generate_descriptive_titles', True):
            return url
        
        # Extract domain and path components
        parsed_url = urlparse(url)
        domain = parsed_url.netloc.replace('www.', '')
        
        path_parts = [p for p in parsed_url.path.split('/') if p]
        
        # Format title
        if not path_parts:
            return domain
        
        # Use last meaningful path component as title
        slug = path_parts[-1].replace('-', ' ').replace('_', ' ')
        
        # Clean up the slug
        slug = re.sub(r'\.\w+$', '', slug)  # Remove file extensions
        
        if not slug:
            return domain
        
        return f"{slug.title()} - {domain}"
    
    def process_file(self, input_file: str, output_dir: str = "output") -> Dict[str, Any]:
        """
        Process input file to extract links and generate QR codes.
        
        Args:
            input_file: Path to input file
            output_dir: Base directory for outputs
            
        Returns:
            Dictionary with processing results
        """
        print(f"Processing file: {input_file}")
        
        # Extract base filename without extension
        base_filename = os.path.splitext(os.path.basename(input_file))[0]
        
        # Create output directory
        file_output_dir = os.path.join(output_dir, base_filename)
        os.makedirs(file_output_dir, exist_ok=True)
        
        # Extract links
        links = self.link_extractor.extract_from_file(input_file)
        print(f"Found {len(links)} unique links")
        
        # Process each link
        qr_data = []
        for url in links:
            # Create sanitized directory name
            dir_name = self.sanitize_filename(url)
            url_dir = os.path.join(file_output_dir, dir_name)
            
            # Generate QR code
            files = self.qr_generator.generate_qr_code(url, url_dir)
            
            # Generate title
            title = self.generate_title(url)
            
            qr_data.append({
                'url': url,
                'title': title,
                'directory': url_dir,
                'files': files
            })
        
        # Generate outputs if configured and links exist
        results = {
            'input_file': input_file,
            'output_dir': file_output_dir,
            'links_found': len(links),
            'qr_data': qr_data,
            'outputs': {}
        }
        
        if qr_data:
            if self.config.get('output_pptx', True):
                pptx_path = self.output_generator.create_powerpoint(qr_data, file_output_dir)
                if pptx_path:
                    results['outputs']['pptx'] = pptx_path
            
            if self.config.get('output_pdf', True):
                pdf_path = self.output_generator.create_pdf(qr_data, file_output_dir)
                if pdf_path:
                    results['outputs']['pdf'] = pdf_path
        
        return results


def main():
    """Main entry point for the command line interface."""
    parser = argparse.ArgumentParser(
        description="Extract links from files and generate QR codes and presentations."
    )
    parser.add_argument(
        "input_file", 
        help="Path to input file (txt, md, csv, pdf, etc.)"
    )
    parser.add_argument(
        "-o", "--output-dir", 
        default="output",
        help="Base directory for outputs (default: output)"
    )
    parser.add_argument(
        "-c", "--config", 
        help="Path to YAML configuration file"
    )
    
    args = parser.parse_args()
    
    # Check if input file exists
    if not os.path.exists(args.input_file):
        print(f"Error: Input file '{args.input_file}' not found.")
        return 1
    
    # Create app instance
    app = LinkProcessorApp(args.config)
    
    # Process file
    result = app.process_file(args.input_file, args.output_dir)
    
    # Print summary
    print("\nProcessing complete!")
    print(f"Processed file: {result['input_file']}")
    print(f"Found {result['links_found']} unique links")
    print(f"Created QR codes in: {result['output_dir']}")
    
    for output_type, path in result['outputs'].items():
        print(f"Generated {output_type.upper()}: {path}")
    
    return 0


if __name__ == "__main__":
    main()